#!/usr/bin/env python3
# 8장 PNG(2x)를 16:9 PPTX 풀블리드 이미지 슬라이드로 조립.
import os
from pptx import Presentation
from pptx.util import Inches

HERE = os.path.dirname(os.path.abspath(__file__))
PNG = os.path.join(HERE, "pptx_png")
N = 10

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

missing = []
added = 0
for i in range(1, N + 1):
    p = os.path.join(PNG, f"slide-{i}.png")
    if not os.path.exists(p):
        missing.append(i)
        continue
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(p, 0, 0, width=prs.slide_width, height=prs.slide_height)
    added += 1

out = os.path.join(HERE, "LMS 외부서비스 결정.pptx")
prs.save(out)
print("saved", out, "| slides:", added)
if missing:
    print("MISSING PNG:", missing)
